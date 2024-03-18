from pptx import Presentation



def extract_slide_titles(pptx_file_path):
    # Create a Presentation object
    prs = Presentation(pptx_file_path)
    
    # Dictionary to store slide titles by index
    slide_titles = {}
    previous_title = None
    
    # Iterate over all slides
    for i, slide in enumerate(prs.slides):
        # Check if the slide has a title
        if slide.shapes.title is not None:
            # Get the slide title
            title = slide.shapes.title.text
            slide_titles[i] = title
            previous_title = title
        else:
            # If the slide does not have a title, use the previous title
            if previous_title is not None:
                slide_titles[i] = previous_title
            else:
                # If there's no previous title, use an empty string
                slide_titles[i] = ""
    
    return slide_titles



# # Example usage:
# pptx_file_path = 'example.pptx'
# slide_titles = extract_slide_titles(pptx_file_path)
# print(slide_titles)
