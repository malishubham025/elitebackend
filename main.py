from flask import Flask, request, jsonify,make_response
import os
import convertapi
from werkzeug.utils import secure_filename
import uuid
from pptx import Presentation
from extracttitles import extract_slide_titles
from text_from_title import generate_explanations
from file_paths import sort_files_and_get_paths
from gemini import generate_image_explanations
from pptx import Presentation
import google.generativeai as genai
from voice import generate_and_save_audio_files
from topic_name import get_topic_name
GOOGLE_API_KEY=os.environ.get('gemini')
genai.configure(api_key=GOOGLE_API_KEY)
model = genai.GenerativeModel('gemini-pro-vision')
# import inx

app = Flask(__name__)

UPLOAD_FOLDER = 'uploads'
SCREENSHOT_FOLDER = 'screenshots'
TEXT_FOLDER = 'text_content'
audio='audio'
if not os.path.exists(UPLOAD_FOLDER):
    os.makedirs(UPLOAD_FOLDER)

if not os.path.exists(SCREENSHOT_FOLDER):
    os.makedirs(SCREENSHOT_FOLDER)

if not os.path.exists(TEXT_FOLDER):
    os.makedirs(TEXT_FOLDER)

convertapi.api_secret = os.environ.get('ppt')

def read_pptx(ppt_file):
    T = ""
    # Function to extract text from slides
    def extract_text_from_slide(slide):
        text = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        text.append(run.text)
        return ' '.join(text)

    # Read the PowerPoint file
    prs = Presentation(ppt_file)

    # Dictionary to store text content by slide numbers
    slide_contents = {}

    # Extract text from each slide and write to a text file
    # with open('textfile.txt', 'w') as f:
    for idx, slide in enumerate(prs.slides):
        slide_text = extract_text_from_slide(slide)
        slide_contents[idx] = slide_text
        # f.write(slide_text + '\n')
        T += slide_text

    return T,slide_contents

def remove_empty_slides(pptx_file_path):
    # Load the PowerPoint presentation
    prs = Presentation(pptx_file_path)
    
    # List to store indices of empty slides
    empty_slide_indices = []
    
    # Iterate over slides and identify empty slides
    for i, slide in enumerate(prs.slides):
        # Check if the slide has no shapes
        if len(slide.shapes) == 0:
            empty_slide_indices.append(i)
    
    # Iterate over the indices of empty slides in reverse order and delete them
    for index in reversed(empty_slide_indices):
        rId = prs.slides._sldIdLst[index].rId
        prs.part.drop_rel(rId)
        del prs.slides._sldIdLst[index]
    
    # Save the modified presentation, overwriting the original file
    prs.save(pptx_file_path)

import shutil
def take_screenshot(file_path, screenshot_folder):
    # Delete the contents of the screenshot folder
    if os.path.exists(screenshot_folder):
        shutil.rmtree(screenshot_folder)
        
    if os.path.exists(audio):
        shutil.rmtree(audio)
    
    os.makedirs(audio)
    # Create the screenshot folder if it doesn't exist
    os.makedirs(screenshot_folder)
    remove_empty_slides(file_path)
    
    print("Taking screenshots")
    convertapi.api_secret = os.environ.get('ppt')
    convertapi.convert('jpg', {
        'File': file_path,
    }, from_format='ppt').save_files(screenshot_folder)


def check_images_presence(pptx_file):
    prs = Presentation(pptx_file)
    images_present = []

    for slide in prs.slides:
        image_found = False
        for shape in slide.shapes:
            if hasattr(shape, "image"):
                image_found = True
                break
        images_present.append(image_found)

    return images_present
from openai import OpenAI



@app.after_request
async def add_cors_headers(response):
    response.headers.add('Access-Control-Allow-Origin', '*')
    response.headers.add('Access-Control-Allow-Methods', 'POST')
    response.headers.add('Access-Control-Allow-Headers', 'Content-Type')
    return response

from flask import jsonify

@app.route('/upload', methods=['POST'])
async def upload_file():
    if 'file' not in request.files:
        return jsonify({'error': 'No file part'}), 400

    file = request.files['file']
    if file.filename == '':
        return jsonify({'error': 'No selected file'}), 400

    if file and file.filename.endswith(('ppt', 'pptx')):
        filename = secure_filename(file.filename)
        file_path = os.path.join(UPLOAD_FOLDER, filename)
        file.save(file_path)

        # Take a screenshot of the uploaded PowerPoint file
        take_screenshot(file_path, SCREENSHOT_FOLDER)
        print("=========================ScreenShot complete============================================")
        images_presence = check_images_presence(file_path)
        print("=======================image presence  complete==========================================")
        print(images_presence) 
        response_data = {
            'message': 'File uploaded successfully',
            'filename': filename,
            'images_presence': images_presence,  # Include the images_presence array in the response
            'file_path':file_path
        }
        response = make_response(jsonify(response_data))
        
        response.set_cookie('file_path', file_path)
        
        return response, 200
    else:
        return jsonify({'error': 'error while processing'}), 400


from flask import request
from explain import  explainations
@app.route('/image_presence', methods=['POST'])
def check_image_presence():
    try:
        request_data = request.json  # Parse JSON data from request body
        images_presence = request_data.get('images_presence')
        file_path = request_data.get('file_path')
        
        # Check if required data is present
        if not images_presence or not file_path:
            return jsonify({'error': 'Missing required data'}), 400

        # Perform operations on the data
        try:
            T,slide_contents = read_pptx(file_path)
            print("=======================Text extraction Complete==========================================")

            Title = extract_slide_titles(file_path)
            print("=======================Title extraction complete==========================================")
            text = ""
            for i in Title:
                text += Title[i] + ","
            topic_name = get_topic_name(T, text)
            print("=======================Topic name==========================================")
            print(topic_name)
            #generate_explanations(slide_contents,topic_name)
            explanations =generate_explanations(slide_contents,topic_name)
            # print(explanations)
            print("=======================Explanations complete==========================================")

            file_paths = sort_files_and_get_paths("screenshots")

            print("============================================================================")
            print(len(explanations))
            print(len(images_presence))
            print(len(Title))
            print(len(file_paths))
            print(topic_name)
            # Prepare response data
            response_data = {
                'message': 'File uploaded successfully',
                'explanations': explanations,
                'images_presence': images_presence,
                'Title': Title,
                'file_paths': file_paths,
                'topic_name':topic_name
            }
            return jsonify(response_data), 200
        except Exception as e:
            return jsonify({'error': str(e)}), 500  # Handle unexpected errors
    except Exception as e:
        return jsonify({'error': str(e)}), 400  # Handle parsing JSON error


@app.route('/image_explain', methods=['POST'])
def image_explain():
    try:
        request_data = request.json  # Parse JSON data from request body
        
        # Check if all required data is present in the request
        required_keys = ['Title', 'images_presence', 'file_paths', 'explanations']
        for key in required_keys:
            if key not in request_data:
                return jsonify({'error': f'Missing required key: {key}'}), 400

        # Extract data from the request
        Title = request_data['Title']
        images_presence = request_data['images_presence']
        file_paths = request_data['file_paths']
        explanations = request_data['explanations']
        print("printing explainations =================================")
        print(explanations)
        # Perform image explanations
        explanations = generate_image_explanations(Title, images_presence, file_paths, model, explanations)
        print("======================image explanation================================")
        
        # Generate and save audio files
        generate_and_save_audio_files(Title, explanations,file_paths)
        print("======================Audio generated=====================================")   

        # Return success message
        return jsonify({'message': 'Image explanation and audio generation completed'}), 200

    except Exception as e:
        # Handle any unexpected errors
        return jsonify({'error': str(e)}), 500
    
from flask import Flask, send_from_directory
from test import upload_file_to_s3
@app.route('/video', methods=['GET'])
def video():
    # return jsonify('response_data'), 200
    try:
        file_path = 'audio/final_video.mp4'  # replace with the path to your file in your local machine
        bucket_name = 'ldgro'
        region = 'ap-south-1'  # specify your AWS region
        access_key_id =os.environ.get("access_key_id")   # your AWS access key ID
        secret_access_key = os.environ.get("secret_access_key") # your AWS secret access key

        url=upload_file_to_s3(file_path, bucket_name, region, access_key_id, secret_access_key)
        # print(folder_name)
        response_data = {
            'url':url
        }
        return jsonify(response_data), 200
    except FileNotFoundError:
        return 'Video not found', 404

@app.route('/signup', methods=['POST'])
def signup():
    try:
        request_data = request.json
        email=(request_data['email'])
        response = make_response(jsonify({'message': 'Signup successful'}), 200)
        response.set_cookie('user_email', email) 
        return response
    except Exception as e:
        return jsonify({'error': str(e)}), 500
    
    


if __name__ == '__main__':
    app.run(debug=True,port=8000)